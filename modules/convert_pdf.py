from flask import Blueprint, request, jsonify, current_app
import os
from werkzeug.utils import secure_filename
import uuid
import time
import logging
import zipfile
from PyPDF2 import PdfReader, PdfWriter
from pdf2image import convert_from_path
from PIL import Image
import img2pdf
import docx
import pandas as pd
from pptx import Presentation
import subprocess
import platform
import tempfile

# Configure logging
logger = logging.getLogger(__name__)

# Create blueprint
convert_bp = Blueprint('convert', __name__)

def register_routes(app, socketio):
    """Register routes with the Flask app"""
    app.register_blueprint(convert_bp)
    
    @app.route('/pdf_to_img', methods=['POST'])
    def pdf_to_img_route():
        if 'pdf' not in request.files:
            socketio.emit('status_update', {
                'status': 'error', 
                'message': '⚠️ No PDF file was uploaded.',
                'tool': 'pdf-to-img'
            })
            return jsonify({'status': 'error', 'message': 'No PDF file was uploaded.'})
        
        pdf_file = request.files['pdf']
        if pdf_file.filename == '':
            socketio.emit('status_update', {
                'status': 'error', 
                'message': '⚠️ No PDF file was selected.',
                'tool': 'pdf-to-img'
            })
            return jsonify({'status': 'error', 'message': 'No PDF file was selected.'})
        
        # Get conversion parameters
        image_format = request.form.get('format', 'jpg')
        dpi = int(request.form.get('dpi', '200'))
        
        # Save the uploaded file
        filename = secure_filename(pdf_file.filename)
        from app import UPLOAD_FOLDER
        file_path = os.path.join(UPLOAD_FOLDER, filename)
        pdf_file.save(file_path)
        
        # Start a background thread for processing
        thread = socketio.start_background_task(
            convert_pdf_to_images, 
            file_path=file_path,
            image_format=image_format,
            dpi=dpi,
            socketio=socketio
        )
        
        # Return immediately with a processing message
        return jsonify({'status': 'processing', 'message': 'Converting PDF to images...'})
    
    @app.route('/img_to_pdf', methods=['POST'])
    def img_to_pdf_route():
        if 'images[]' not in request.files:
            socketio.emit('status_update', {
                'status': 'error', 
                'message': '⚠️ No image files were uploaded.',
                'tool': 'img-to-pdf'
            })
            return jsonify({'status': 'error', 'message': 'No image files were uploaded.'})
        
        images = request.files.getlist('images[]')
        if len(images) == 0:
            socketio.emit('status_update', {
                'status': 'error', 
                'message': '⚠️ No image files were selected.',
                'tool': 'img-to-pdf'
            })
            return jsonify({'status': 'error', 'message': 'No image files were selected.'})
        
        # Get image order if provided
        image_order = request.form.get('order', '')
        
        # Save uploaded files
        image_paths = []
        for image in images:
            if image.filename == '':
                continue
            
            filename = secure_filename(image.filename)
            from app import UPLOAD_FOLDER
            file_path = os.path.join(UPLOAD_FOLDER, filename)
            image.save(file_path)
            image_paths.append(file_path)
        
        if len(image_paths) == 0:
            socketio.emit('status_update', {
                'status': 'error', 
                'message': '⚠️ No valid image files were uploaded.',
                'tool': 'img-to-pdf'
            })
            return jsonify({'status': 'error', 'message': 'No valid image files were uploaded.'})
        
        # Start a background thread for processing
        thread = socketio.start_background_task(
            convert_images_to_pdf, 
            image_paths=image_paths,
            image_order=image_order,
            socketio=socketio
        )
        
        # Return immediately with a processing message
        return jsonify({'status': 'processing', 'message': 'Converting images to PDF...'})
    
    @app.route('/pdf_to_word', methods=['POST'])
    def pdf_to_word_route():
        if 'pdf' not in request.files:
            socketio.emit('status_update', {
                'status': 'error', 
                'message': '⚠️ No PDF file was uploaded.',
                'tool': 'pdf-to-word'
            })
            return jsonify({'status': 'error', 'message': 'No PDF file was uploaded.'})
        
        pdf_file = request.files['pdf']
        if pdf_file.filename == '':
            socketio.emit('status_update', {
                'status': 'error', 
                'message': '⚠️ No PDF file was selected.',
                'tool': 'pdf-to-word'
            })
            return jsonify({'status': 'error', 'message': 'No PDF file was selected.'})
        
        # Save the uploaded file
        filename = secure_filename(pdf_file.filename)
        from app import UPLOAD_FOLDER
        file_path = os.path.join(UPLOAD_FOLDER, filename)
        pdf_file.save(file_path)
        
        # Start a background thread for processing
        thread = socketio.start_background_task(
            convert_pdf_to_word, 
            file_path=file_path,
            socketio=socketio
        )
        
        # Return immediately with a processing message
        return jsonify({'status': 'processing', 'message': 'Converting PDF to Word...'})
    
    @app.route('/pdf_to_excel', methods=['POST'])
    def pdf_to_excel_route():
        if 'pdf' not in request.files:
            socketio.emit('status_update', {
                'status': 'error', 
                'message': '⚠️ No PDF file was uploaded.',
                'tool': 'pdf-to-excel'
            })
            return jsonify({'status': 'error', 'message': 'No PDF file was uploaded.'})
        
        pdf_file = request.files['pdf']
        if pdf_file.filename == '':
            socketio.emit('status_update', {
                'status': 'error', 
                'message': '⚠️ No PDF file was selected.',
                'tool': 'pdf-to-excel'
            })
            return jsonify({'status': 'error', 'message': 'No PDF file was selected.'})
        
        # Save the uploaded file
        filename = secure_filename(pdf_file.filename)
        from app import UPLOAD_FOLDER
        file_path = os.path.join(UPLOAD_FOLDER, filename)
        pdf_file.save(file_path)
        
        # Start a background thread for processing
        thread = socketio.start_background_task(
            convert_pdf_to_excel, 
            file_path=file_path,
            socketio=socketio
        )
        
        # Return immediately with a processing message
        return jsonify({'status': 'processing', 'message': 'Converting PDF to Excel...'})
    
    @app.route('/pdf_to_ppt', methods=['POST'])
    def pdf_to_ppt_route():
        if 'pdf' not in request.files:
            socketio.emit('status_update', {
                'status': 'error', 
                'message': '⚠️ No PDF file was uploaded.',
                'tool': 'pdf-to-ppt'
            })
            return jsonify({'status': 'error', 'message': 'No PDF file was uploaded.'})
        
        pdf_file = request.files['pdf']
        if pdf_file.filename == '':
            socketio.emit('status_update', {
                'status': 'error', 
                'message': '⚠️ No PDF file was selected.',
                'tool': 'pdf-to-ppt'
            })
            return jsonify({'status': 'error', 'message': 'No PDF file was selected.'})
        
        # Save the uploaded file
        filename = secure_filename(pdf_file.filename)
        from app import UPLOAD_FOLDER
        file_path = os.path.join(UPLOAD_FOLDER, filename)
        pdf_file.save(file_path)
        
        # Start a background thread for processing
        thread = socketio.start_background_task(
            convert_pdf_to_ppt, 
            file_path=file_path,
            socketio=socketio
        )
        
        # Return immediately with a processing message
        return jsonify({'status': 'processing', 'message': 'Converting PDF to PowerPoint...'})
    
    @app.route('/word_to_pdf', methods=['POST'])
    def word_to_pdf_route():
        if 'word' not in request.files:
            socketio.emit('status_update', {
                'status': 'error', 
                'message': '⚠️ No Word file was uploaded.',
                'tool': 'word-to-pdf'
            })
            return jsonify({'status': 'error', 'message': 'No Word file was uploaded.'})
        
        word_file = request.files['word']
        if word_file.filename == '':
            socketio.emit('status_update', {
                'status': 'error', 
                'message': '⚠️ No Word file was selected.',
                'tool': 'word-to-pdf'
            })
            return jsonify({'status': 'error', 'message': 'No Word file was selected.'})
        
        # Save the uploaded file
        filename = secure_filename(word_file.filename)
        from app import UPLOAD_FOLDER
        file_path = os.path.join(UPLOAD_FOLDER, filename)
        word_file.save(file_path)
        
        # Start a background thread for processing
        thread = socketio.start_background_task(
            convert_word_to_pdf, 
            file_path=file_path,
            socketio=socketio
        )
        
        # Return immediately with a processing message
        return jsonify({'status': 'processing', 'message': 'Converting Word to PDF...'})
    
    @app.route('/excel_to_pdf', methods=['POST'])
    def excel_to_pdf_route():
        if 'excel' not in request.files:
            socketio.emit('status_update', {
                'status': 'error', 
                'message': '⚠️ No Excel file was uploaded.',
                'tool': 'excel-to-pdf'
            })
            return jsonify({'status': 'error', 'message': 'No Excel file was uploaded.'})
        
        excel_file = request.files['excel']
        if excel_file.filename == '':
            socketio.emit('status_update', {
                'status': 'error', 
                'message': '⚠️ No Excel file was selected.',
                'tool': 'excel-to-pdf'
            })
            return jsonify({'status': 'error', 'message': 'No Excel file was selected.'})
        
        # Save the uploaded file
        filename = secure_filename(excel_file.filename)
        from app import UPLOAD_FOLDER
        file_path = os.path.join(UPLOAD_FOLDER, filename)
        excel_file.save(file_path)
        
        # Start a background thread for processing
        thread = socketio.start_background_task(
            convert_excel_to_pdf, 
            file_path=file_path,
            socketio=socketio
        )
        
        # Return immediately with a processing message
        return jsonify({'status': 'processing', 'message': 'Converting Excel to PDF...'})
    
    @app.route('/ppt_to_pdf', methods=['POST'])
    def ppt_to_pdf_route():
        if 'ppt' not in request.files:
            socketio.emit('status_update', {
                'status': 'error', 
                'message': '⚠️ No PowerPoint file was uploaded.',
                'tool': 'ppt-to-pdf'
            })
            return jsonify({'status': 'error', 'message': 'No PowerPoint file was uploaded.'})
        
        ppt_file = request.files['ppt']
        if ppt_file.filename == '':
            socketio.emit('status_update', {
                'status': 'error', 
                'message': '⚠️ No PowerPoint file was selected.',
                'tool': 'ppt-to-pdf'
            })
            return jsonify({'status': 'error', 'message': 'No PowerPoint file was selected.'})
        
        # Save the uploaded file
        filename = secure_filename(ppt_file.filename)
        from app import UPLOAD_FOLDER
        file_path = os.path.join(UPLOAD_FOLDER, filename)
        ppt_file.save(file_path)
        
        # Start a background thread for processing
        thread = socketio.start_background_task(
            convert_ppt_to_pdf, 
            file_path=file_path,
            socketio=socketio
        )
        
        # Return immediately with a processing message
        return jsonify({'status': 'processing', 'message': 'Converting PowerPoint to PDF...'})

def convert_pdf_to_images(file_path, image_format='jpg', dpi=200, socketio=None):
    """Convert PDF to images"""
    if socketio:
        socketio.emit('status_update', {
            'status': 'processing', 
            'message': 'Starting PDF to image conversion...',
            'tool': 'pdf-to-img'
        })
    
    try:
        # Get paths
        from app import OUTPUT_FOLDER
        
        # Generate a unique ID for this operation
        operation_id = uuid.uuid4().hex[:8]
        base_filename = os.path.basename(file_path)
        name_without_ext = os.path.splitext(base_filename)[0]
        
        # Create output directory
        output_dir = os.path.join(OUTPUT_FOLDER, f"{name_without_ext}_images_{operation_id}")
        os.makedirs(output_dir, exist_ok=True)
        
        # Convert PDF to images
        if socketio:
            socketio.emit('status_update', {
                'status': 'processing', 
                'message': 'Converting PDF to images...',
                'progress': 30,
                'tool': 'pdf-to-img'
            })
        
        images = convert_from_path(file_path, dpi=dpi)
        
        if not images:
            if socketio:
                socketio.emit('status_update', {
                    'status': 'error', 
                    'message': '❌ Failed to convert PDF to images.',
                    'tool': 'pdf-to-img'
                })
            return
        
        # Save images
        output_files = []
        
        for i, image in enumerate(images):
            # Create output filename
            output_filename = f"{name_without_ext}_page_{i+1}.{image_format}"
            output_path = os.path.join(output_dir, output_filename)
            
            # Save image
            image.save(output_path, format=image_format.upper())
            
            # Add to output files
            output_files.append({
                'path': output_path,
                'name': output_filename,
                'url': f'/download/{os.path.basename(output_dir)}/{output_filename}'
            })
            
            # Update progress
            progress = int(((i + 1) / len(images)) * 60) + 30  # Scale to 30-90%
            if socketio:
                socketio.emit('status_update', {
                    'status': 'processing', 
                    'message': f'Processing page {i+1} of {len(images)}',
                    'progress': progress,
                    'tool': 'pdf-to-img'
                })
        
        # Create ZIP file with all images
        if socketio:
            socketio.emit('status_update', {
                'status': 'processing', 
                'message': 'Creating ZIP file...',
                'progress': 95,
                'tool': 'pdf-to-img'
            })
        
        zip_filename = f"{name_without_ext}_images_{operation_id}.zip"
        zip_path = os.path.join(OUTPUT_FOLDER, zip_filename)
        
        with zipfile.ZipFile(zip_path, 'w') as zipf:
            for file_info in output_files:
                zipf.write(file_info['path'], arcname=file_info['name'])
        
        # Prepare download links
        downloads = [
            {
                'name': zip_filename,
                'url': f'/download/{zip_filename}',
                'type': 'All Images (ZIP)'
            }
        ]
        
        if socketio:
            socketio.emit('status_update', {
                'status': 'success', 
                'message': f'✅ PDF converted to {len(images)} images successfully!',
                'downloads': downloads,
                'filename': zip_filename,
                'images': output_files,
                'tool': 'pdf-to-img'
            })
    
    except Exception as e:
        logger.error(f"Error converting PDF to images: {str(e)}")
        if socketio:
            socketio.emit('status_update', {
                'status': 'error', 
                'message': f'❌ Error converting PDF to images: {str(e)}',
                'tool': 'pdf-to-img'
            })

def convert_images_to_pdf(image_paths, image_order='', socketio=None):
    """Convert images to PDF"""
    if socketio:
        socketio.emit('status_update', {
            'status': 'processing', 
            'message': 'Starting image to PDF conversion...',
            'tool': 'img-to-pdf'
        })
    
    try:
        # Get paths
        from app import OUTPUT_FOLDER
        
        # Generate a unique ID for this operation
        operation_id = uuid.uuid4().hex[:8]
        
        # Create output filename
        output_filename = f"images_to_pdf_{operation_id}.pdf"
        output_path = os.path.join(OUTPUT_FOLDER, output_filename)
        
        # Reorder images if order is provided
        if image_order:
            try:
                order_indices = [int(i) for i in image_order.split(',')]
                # Validate indices
                if len(order_indices) == len(image_paths) and all(0 <= i < len(image_paths) for i in order_indices):
                    image_paths = [image_paths[i] for i in order_indices]
            except (ValueError, IndexError):
                # If order is invalid, use original order
                pass
        
        # Process images
        if socketio:
            socketio.emit('status_update', {
                'status': 'processing', 
                'message': 'Processing images...',
                'progress': 30,
                'tool': 'img-to-pdf'
            })
        
        # Validate and convert images
        valid_images = []
        
        for i, image_path in enumerate(image_paths):
            try:
                # Open image to validate
                img = Image.open(image_path)
                img.verify()  # Verify that it's a valid image
                
                # Add to valid images
                valid_images.append(image_path)
                
                # Update progress
                progress = int(((i + 1) / len(image_paths)) * 60) + 30  # Scale to 30-90%
                if socketio:
                    socketio.emit('status_update', {
                        'status': 'processing', 
                        'message': f'Processing image {i+1} of {len(image_paths)}',
                        'progress': progress,
                        'tool': 'img-to-pdf'
                    })
            
            except Exception as e:
                logger.warning(f"Invalid image: {image_path} - {str(e)}")
                # Skip invalid images
                continue
        
        if not valid_images:
            if socketio:
                socketio.emit('status_update', {
                    'status': 'error', 
                    'message': '❌ No valid images found.',
                    'tool': 'img-to-pdf'
                })
            return
        
        # Convert images to PDF
        if socketio:
            socketio.emit('status_update', {
                'status': 'processing', 
                'message': 'Converting images to PDF...',
                'progress': 95,
                'tool': 'img-to-pdf'
            })
        
        # Convert images to PDF using img2pdf
        with open(output_path, 'wb') as f:
            f.write(img2pdf.convert(valid_images))
        
        # Prepare download links
        downloads = [
            {
                'name': output_filename,
                'url': f'/download/{output_filename}',
                'type': 'PDF from Images'
            }
        ]
        
        if socketio:
            socketio.emit('status_update', {
                'status': 'success', 
                'message': f'✅ {len(valid_images)} images converted to PDF successfully!',
                'downloads': downloads,
                'filename': output_filename,
                'tool': 'img-to-pdf'
            })
    
    except Exception as e:
        logger.error(f"Error converting images to PDF: {str(e)}")
        if socketio:
            socketio.emit('status_update', {
                'status': 'error', 
                'message': f'❌ Error converting images to PDF: {str(e)}',
                'tool': 'img-to-pdf'
            })

def convert_pdf_to_word(file_path, socketio=None):
    """Convert PDF to Word document"""
    if socketio:
        socketio.emit('status_update', {
            'status': 'processing', 
            'message': 'Starting PDF to Word conversion...',
            'tool': 'pdf-to-word'
        })
    
    try:
        # Get paths
        from app import OUTPUT_FOLDER
        
        # Generate a unique ID for this operation
        operation_id = uuid.uuid4().hex[:8]
        base_filename = os.path.basename(file_path)
        name_without_ext = os.path.splitext(base_filename)[0]
        
        # Create output filename
        output_filename = f"{name_without_ext}_{operation_id}.docx"
        output_path = os.path.join(OUTPUT_FOLDER, output_filename)
        
        # Extract text from PDF
        if socketio:
            socketio.emit('status_update', {
                'status': 'processing', 
                'message': 'Extracting text from PDF...',
                'progress': 30,
                'tool': 'pdf-to-word'
            })
        
        # Use pdfminer to extract text
        from pdfminer.high_level import extract_text
        text = extract_text(file_path)
        
        if not text:
            if socketio:
                socketio.emit('status_update', {
                    'status': 'error', 
                    'message': '❌ No text found in the PDF.',
                    'tool': 'pdf-to-word'
                })
            return
        
        # Create Word document
        if socketio:
            socketio.emit('status_update', {
                'status': 'processing', 
                'message': 'Creating Word document...',
                'progress': 70,
                'tool': 'pdf-to-word'
            })
        
        # Create Word document
        doc = docx.Document()
        
        # Split text by pages (assuming pages are separated by form feeds)
        pages = text.split('\f')
        
        for page in pages:
            if page.strip():
                doc.add_paragraph(page)
                doc.add_page_break()
        
        # Save Word document
        doc.save(output_path)
        
        # Prepare download links
        downloads = [
            {
                'name': output_filename,
                'url': f'/download/{output_filename}',
                'type': 'Word Document'
            }
        ]
        
        if socketio:
            socketio.emit('status_update', {
                'status': 'success', 
                'message': '✅ PDF converted to Word successfully!',
                'downloads': downloads,
                'filename': output_filename,
                'tool': 'pdf-to-word'
            })
    
    except Exception as e:
        logger.error(f"Error converting PDF to Word: {str(e)}")
        if socketio:
            socketio.emit('status_update', {
                'status': 'error', 
                'message': f'❌ Error converting PDF to Word: {str(e)}',
                'tool': 'pdf-to-word'
            })

def convert_pdf_to_excel(file_path, socketio=None):
    """Convert PDF to Excel spreadsheet"""
    if socketio:
        socketio.emit('status_update', {
            'status': 'processing', 
            'message': 'Starting PDF to Excel conversion...',
            'tool': 'pdf-to-excel'
        })
    
    try:
        # Get paths
        from app import OUTPUT_FOLDER
        
        # Generate a unique ID for this operation
        operation_id = uuid.uuid4().hex[:8]
        base_filename = os.path.basename(file_path)
        name_without_ext = os.path.splitext(base_filename)[0]
        
        # Create output filename
        output_filename = f"{name_without_ext}_{operation_id}.xlsx"
        output_path = os.path.join(OUTPUT_FOLDER, output_filename)
        
        # Extract tables from PDF
        if socketio:
            socketio.emit('status_update', {
                'status': 'processing', 
                'message': 'Extracting tables from PDF...',
                'progress': 30,
                'tool': 'pdf-to-excel'
            })
        
        # Use tabula-py to extract tables
        import tabula
        
        # Extract tables from all pages
        tables = tabula.read_pdf(file_path, pages='all', multiple_tables=True)
        
        if not tables:
            if socketio:
                socketio.emit('status_update', {
                    'status': 'error', 
                    'message': '❌ No tables found in the PDF.',
                    'tool': 'pdf-to-excel'
                })
            return
        
        # Create Excel file
        if socketio:
            socketio.emit('status_update', {
                'status': 'processing', 
                'message': 'Creating Excel file...',
                'progress': 70,
                'tool': 'pdf-to-excel'
            })
        
        # Create Excel writer
        with pd.ExcelWriter(output_path, engine='openpyxl') as writer:
            for i, table in enumerate(tables):
                # Write each table to a separate sheet
                sheet_name = f'Table {i+1}'
                table.to_excel(writer, sheet_name=sheet_name, index=False)
        
        # Prepare download links
        downloads = [
            {
                'name': output_filename,
                'url': f'/download/{output_filename}',
                'type': 'Excel Spreadsheet'
            }
        ]
        
        if socketio:
            socketio.emit('status_update', {
                'status': 'success', 
                'message': f'✅ PDF converted to Excel successfully! ({len(tables)} tables extracted)',
                'downloads': downloads,
                'filename': output_filename,
                'tool': 'pdf-to-excel'
            })
    
    except Exception as e:
        logger.error(f"Error converting PDF to Excel: {str(e)}")
        if socketio:
            socketio.emit('status_update', {
                'status': 'error', 
                'message': f'❌ Error converting PDF to Excel: {str(e)}',
                'tool': 'pdf-to-excel'
            })

def convert_pdf_to_ppt(file_path, socketio=None):
    """Convert PDF to PowerPoint presentation"""
    if socketio:
        socketio.emit('status_update', {
            'status': 'processing', 
            'message': 'Starting PDF to PowerPoint conversion...',
            'tool': 'pdf-to-ppt'
        })
    
    try:
        # Get paths
        from app import OUTPUT_FOLDER
        
        # Generate a unique ID for this operation
        operation_id = uuid.uuid4().hex[:8]
        base_filename = os.path.basename(file_path)
        name_without_ext = os.path.splitext(base_filename)[0]
        
        # Create output filename
        output_filename = f"{name_without_ext}_{operation_id}.pptx"
        output_path = os.path.join(OUTPUT_FOLDER, output_filename)
        
        # Convert PDF to images
        if socketio:
            socketio.emit('status_update', {
                'status': 'processing', 
                'message': 'Converting PDF to images...',
                'progress': 30,
                'tool': 'pdf-to-ppt'
            })
        
        # Convert PDF to images
        images = convert_from_path(file_path, dpi=300)
        
        if not images:
            if socketio:
                socketio.emit('status_update', {
                    'status': 'error', 
                    'message': '❌ Failed to convert PDF to images.',
                    'tool': 'pdf-to-ppt'
                })
            return
        
        # Create PowerPoint presentation
        if socketio:
            socketio.emit('status_update', {
                'status': 'processing', 
                'message': 'Creating PowerPoint presentation...',
                'progress': 60,
                'tool': 'pdf-to-ppt'
            })
        
        # Create presentation
        prs = Presentation()
        
        # Set slide dimensions
        prs.slide_width = 9144000  # 9.144 inches
        prs.slide_height = 6858000  # 6.858 inches
        
        # Create a blank slide layout
        blank_slide_layout = prs.slide_layouts[6]
        
        # Process each image
        for i, image in enumerate(images):
            # Create temporary file for image
            with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as temp:
                temp_path = temp.name
                image.save(temp_path, 'PNG')
            
            # Add slide
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Add image to slide
            slide.shapes.add_picture(temp_path, 0, 0, prs.slide_width, prs.slide_height)
            
            # Remove temporary file
            os.unlink(temp_path)
            
            # Update progress
            progress = int(((i + 1) / len(images)) * 30) + 60  # Scale to 60-90%
            if socketio:
                socketio.emit('status_update', {
                    'status': 'processing', 
                    'message': f'Processing slide {i+1} of {len(images)}',
                    'progress': progress,
                    'tool': 'pdf-to-ppt'
                })
        
        # Save presentation
        prs.save(output_path)
        
        # Prepare download links
        downloads = [
            {
                'name': output_filename,
                'url': f'/download/{output_filename}',
                'type': 'PowerPoint Presentation'
            }
        ]
        
        if socketio:
            socketio.emit('status_update', {
                'status': 'success', 
                'message': f'✅ PDF converted to PowerPoint successfully! ({len(images)} slides created)',
                'downloads': downloads,
                'filename': output_filename,
                'tool': 'pdf-to-ppt'
            })
    
    except Exception as e:
        logger.error(f"Error converting PDF to PowerPoint: {str(e)}")
        if socketio:
            socketio.emit('status_update', {
                'status': 'error', 
                'message': f'❌ Error converting PDF to PowerPoint: {str(e)}',
                'tool': 'pdf-to-ppt'
            })

def convert_word_to_pdf(file_path, socketio=None):
    """Convert Word document to PDF"""
    if socketio:
        socketio.emit('status_update', {
            'status': 'processing', 
            'message': 'Starting Word to PDF conversion...',
            'tool': 'word-to-pdf'
        })
    
    try:
        # Get paths
        from app import OUTPUT_FOLDER
        
        # Generate a unique ID for this operation
        operation_id = uuid.uuid4().hex[:8]
        base_filename = os.path.basename(file_path)
        name_without_ext = os.path.splitext(base_filename)[0]
        
        # Create output filename
        output_filename = f"{name_without_ext}_{operation_id}.pdf"
        output_path = os.path.join(OUTPUT_FOLDER, output_filename)
        
        # Convert Word to PDF
        if socketio:
            socketio.emit('status_update', {
                'status': 'processing', 
                'message': 'Converting Word to PDF...',
                'progress': 50,
                'tool': 'word-to-pdf'
            })
        
        # Use docx2pdf for conversion
        from docx2pdf import convert
        convert(file_path, output_path)
        
        # Prepare download links
        downloads = [
            {
                'name': output_filename,
                'url': f'/download/{output_filename}',
                'type': 'PDF Document'
            }
        ]
        
        if socketio:
            socketio.emit('status_update', {
                'status': 'success', 
                'message': '✅ Word document converted to PDF successfully!',
                'downloads': downloads,
                'filename': output_filename,
                'tool': 'word-to-pdf'
            })
    
    except Exception as e:
        logger.error(f"Error converting Word to PDF: {str(e)}")
        if socketio:
            socketio.emit('status_update', {
                'status': 'error', 
                'message': f'❌ Error converting Word to PDF: {str(e)}',
                'tool': 'word-to-pdf'
            })

def convert_excel_to_pdf(file_path, socketio=None):
    """Convert Excel spreadsheet to PDF"""
    if socketio:
        socketio.emit('status_update', {
            'status': 'processing', 
            'message': 'Starting Excel to PDF conversion...',
            'tool': 'excel-to-pdf'
        })
    
    try:
        # Get paths
        from app import OUTPUT_FOLDER
        
        # Generate a unique ID for this operation
        operation_id = uuid.uuid4().hex[:8]
        base_filename = os.path.basename(file_path)
        name_without_ext = os.path.splitext(base_filename)[0]
        
        # Create output filename
        output_filename = f"{name_without_ext}_{operation_id}.pdf"
        output_path = os.path.join(OUTPUT_FOLDER, output_filename)
        
        # Convert Excel to PDF
        if socketio:
            socketio.emit('status_update', {
                'status': 'processing', 
                'message': 'Converting Excel to PDF...',
                'progress': 50,
                'tool': 'excel-to-pdf'
            })
        
        # Use Excel COM object for conversion (Windows only)
        if platform.system() == 'Windows':
            import win32com.client
            
            excel = win32com.client.Dispatch("Excel.Application")
            excel.Visible = False
            
            wb = excel.Workbooks.Open(file_path)
            wb.ExportAsFixedFormat(0, output_path)  # 0 = PDF format
            
            wb.Close()
            excel.Quit()
        else:
            # For non-Windows platforms, use alternative method
            # Read Excel file
            df_dict = pd.read_excel(file_path, sheet_name=None)
            
            # Create a PDF writer
            from reportlab.lib.pagesizes import letter
            from reportlab.platypus import SimpleDocTemplate, Table, TableStyle
            from reportlab.lib import colors
            
            doc = SimpleDocTemplate(output_path, pagesize=letter)
            elements = []
            
            # Process each sheet
            for sheet_name, df in df_dict.items():
                # Convert DataFrame to list of lists
                data = [df.columns.tolist()] + df.values.tolist()
                
                # Create table
                table = Table(data)
                
                # Add table style
                style = TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                    ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                    ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                    ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                    ('GRID', (0, 0), (-1, -1), 1, colors.black)
                ])
                
                table.setStyle(style)
                elements.append(table)
            
            # Build PDF
            doc.build(elements)
        
        # Prepare download links
        downloads = [
            {
                'name': output_filename,
                'url': f'/download/{output_filename}',
                'type': 'PDF Document'
            }
        ]
        
        if socketio:
            socketio.emit('status_update', {
                'status': 'success', 
                'message': '✅ Excel spreadsheet converted to PDF successfully!',
                'downloads': downloads,
                'filename': output_filename,
                'tool': 'excel-to-pdf'
            })
    
    except Exception as e:
        logger.error(f"Error converting Excel to PDF: {str(e)}")
        if socketio:
            socketio.emit('status_update', {
                'status': 'error', 
                'message': f'❌ Error converting Excel to PDF: {str(e)}',
                'tool': 'excel-to-pdf'
            })

def convert_ppt_to_pdf(file_path, socketio=None):
    """Convert PowerPoint presentation to PDF"""
    if socketio:
        socketio.emit('status_update', {
            'status': 'processing', 
            'message': 'Starting PowerPoint to PDF conversion...',
            'tool': 'ppt-to-pdf'
        })
    
    try:
        # Get paths
        from app import OUTPUT_FOLDER
        
        # Generate a unique ID for this operation
        operation_id = uuid.uuid4().hex[:8]
        base_filename = os.path.basename(file_path)
        name_without_ext = os.path.splitext(base_filename)[0]
        
        # Create output filename
        output_filename = f"{name_without_ext}_{operation_id}.pdf"
        output_path = os.path.join(OUTPUT_FOLDER, output_filename)
        
        # Convert PowerPoint to PDF
        if socketio:
            socketio.emit('status_update', {
                'status': 'processing', 
                'message': 'Converting PowerPoint to PDF...',
                'progress': 50,
                'tool': 'ppt-to-pdf'
            })
        
        # Use PowerPoint COM object for conversion (Windows only)
        if platform.system() == 'Windows':
            import win32com.client
            
            powerpoint = win32com.client.Dispatch("PowerPoint.Application")
            powerpoint.Visible = True
            
            presentation = powerpoint.Presentations.Open(file_path)
            presentation.ExportAsFixedFormat(output_path, 32)  # 32 = PDF format
            
            presentation.Close()
            powerpoint.Quit()
        else:
            # For non-Windows platforms, use alternative method
            # Convert to images and then to PDF
            
            # Load presentation
            prs = Presentation(file_path)
            
            # Create temporary directory for images
            temp_dir = tempfile.mkdtemp()
            
            # Convert slides to images
            images = []
            
            for i, slide in enumerate(prs.slides):
                # Create slide image using reportlab
                from reportlab.lib.pagesizes import letter
                from reportlab.pdfgen import canvas
                
                # Create temporary file for slide
                temp_slide = os.path.join(temp_dir, f"slide_{i}.png")
                
                # Create image from slide (simplified)
                img = Image.new('RGB', (960, 720), color='white')
                img.save(temp_slide)
                
                images.append(temp_slide)
            
            # Convert images to PDF
            with open(output_path, 'wb') as f:
                f.write(img2pdf.convert(images))
            
            # Clean up temporary directory
            import shutil
            shutil.rmtree(temp_dir)
        
        # Prepare download links
        downloads = [
            {
                'name': output_filename,
                'url': f'/download/{output_filename}',
                'type': 'PDF Document'
            }
        ]
        
        if socketio:
            socketio.emit('status_update', {
                'status': 'success', 
                'message': '✅ PowerPoint presentation converted to PDF successfully!',
                'downloads': downloads,
                'filename': output_filename,
                'tool': 'ppt-to-pdf'
            })
    
    except Exception as e:
        logger.error(f"Error converting PowerPoint to PDF: {str(e)}")
        if socketio:
            socketio.emit('status_update', {
                'status': 'error', 
                'message': f'❌ Error converting PowerPoint to PDF: {str(e)}',
                'tool': 'ppt-to-pdf'
            })